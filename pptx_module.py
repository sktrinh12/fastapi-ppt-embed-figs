from pptx import Presentation
from pptx.util import Cm, Pt
from PIL import Image
from generate_stats_figs import *

def insert_figure(infile, outfile, dct_props: dict):
    """ Take the input powerpoint file and insert a figure
    into the numbered slide
    """
    figure_img = dct_props['filename']
    left = dct_props['left_spacing']
    top = dct_props['top_spacing']
    width = dct_props['width']
    height = dct_props['height']
    prs = Presentation(infile)
    slide = prs.slides[dct_props['slide_number']]
    pix = slide.shapes.add_picture(figure_img, left, top, width, height)
    prs.save(outfile)
    print(f'Inserted {figure_img}')

def make_dct_props(plot_types_list, cwd, timestamp):
    addn_types = ['regression', 'histogram', 'scatter']
    exportdir = os.path.join(cwd, 'exports')
    dct_props = {pltype: {} for pltype in plot_types_list + addn_types}
    basedir = timestamp
    for pt in dct_props.keys():
        dct_props[pt]['left_spacing'] = Cm(2.5)
        dct_props[pt]['top_spacing'] = Cm(3.75)
        dct_props[pt]['height'] = Cm(10.5)
        dct_props[pt]['width'] = Cm(12)
        dct_props[pt]['slide_number'] = 1
        if pt in ['histogram', 'scatter']:
            exportdir = cwd
            basedir = 'uploads'
        dct_props[pt]['filename'] = os.path.join(exportdir, basedir, f'{pt}.png')
    return dct_props

def make_table_content(stats_dct):
    """
    generate the text that will go into each cell of the ppt table
    """
    table_content = [[None]*3]*3
    table_content[0] = ['% of 4C Ref. MFI Threshold', 'No confidence interval', 'Lower 95% CI']
    table_content[1] = ['75'] + [f'{v:.3g}' for v in calc_thresCI(75, stats_dct)]
    table_content[2] = ['80'] + [f'{v:.3g}' for v in calc_thresCI(80, stats_dct)]
    return table_content

def update_dct_props(dct_props, table_content, rm_si_sn):
    dct_props['mfi_v_time']['left_spacing'] = Cm(20)
    dct_props['pos_v_conc']['slide_number'] = 2
    dct_props['pct_4c_ref']['slide_number'] = 2
    dct_props['pct_4c_ref']['left_spacing'] = Cm(20)
    # print(f'rm_si_sn: {rm_si_sn}')
    if rm_si_sn == False:
        dct_props['stain_index']['slide_number'] = 3
        dct_props['signal_noise']['left_spacing'] = Cm(20)
        dct_props['signal_noise']['slide_number'] = 3
    dct_props['regression']['height'] = None
    dct_props['regression']['width'] = None
    dct_props['regression']['top_spacing'] = Cm(2)
    dct_props['regression']['slide_number'] = 3 if rm_si_sn else 4
    dct_props['regression']['table_cell'] = table_content
    dct_props['histogram']['slide_number'] = 4 if rm_si_sn else 5
    dct_props['histogram']['left_spacing'] = -Cm(0.5)
    dct_props['histogram']['top_spacing'] = Cm(5.75)
    dct_props['histogram']['width'] = Cm(34)
    dct_props['histogram']['height'] = None
    dct_props['scatter']['width'] = Cm(16)
    dct_props['scatter']['height'] = Cm(18)
    dct_props['scatter']['slide_number'] = 5 if rm_si_sn else 6
    dct_props['scatter']['top_spacing'] = Cm(1.1)
    dct_props['scatter']['left_spacing'] = Cm(11)
    return dct_props

def get_img_dims(img_file, wd):
    im=Image.open(os.path.join(wd, img_file))
    return list(im.size)

def purge_folder(wd):
    """
    clear the contents of the export directory prior to renaming
    """
    for filename in os.listdir(wd):
        file_path = os.path.join(wd, filename)
        try:
            if os.path.isfile(file_path):
                os.remove(file_path)
                print(f'removed {filename}')
        except Exception as e:
            print(f'Faield to delete {file_path}. Reason: {e}')

def rename_histo_scatter_imgs(wd):
    """
    rename the histogram and scatter figures to a standard name in order to embed into power point slides
    slides
    """
    height_dim_id = 600
    histo_filename = 'histogram.png'
    scatter_filename = 'scatter.png'
    img_files = [img_fi for img_fi in os.listdir(wd) if img_fi.endswith('.png')]
    assert len(img_files) == 2, "Did not find exactly 2 .png files in the directory"
    img_props = [[img] + get_img_dims(img,wd) for img in img_files]
    histo_img = next((i for i in img_props if i[2] < height_dim_id), None)[0]
    os.rename(os.path.join(wd, histo_img), os.path.join(wd, histo_filename))
    print(f'renamed: {histo_img} -> {histo_filename}')
    scatter_img = next((i for i in img_props if i[0] != histo_img), None)[0]
    os.rename(os.path.join(wd,scatter_img), os.path.join(wd,scatter_filename))
    print(f'renamed: {scatter_img} -> {scatter_filename}')

def insert_regr_table(infile, outfile, dct_props):
    rows, cols = 4, 3
    prs = Presentation(infile)
    slide = prs.slides[dct_props['slide_number']]
    # setting width and height here can be arbitrary so that doesn't throw None error
    left, top, width, height = Cm(24), Cm(7.25), Cm(8), Cm(1)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # merge top row for a title
    top_cell = table.cell(0, 0)
    last_cell = table.cell(0, 2)
    top_cell.merge(last_cell)
    top_cell.text = 'Shelf-life based on:'
    top_cell.text_frame.paragraphs[0].font.size = Pt(14)
    # set column widths more finely here
    table.columns[0].width = Cm(3.25)
    table.columns[1].width = Cm(3)
    table.columns[2].width = Cm(1.85)
    # loop thru other remaining cells to populate
    for i in range(1,rows):
        for j in range(cols):
            cell = table.cell(i, j)
            cell.text = dct_props['table_cell'][i-1][j]
            textfr = cell.text_frame.paragraphs[0].font
            textfr.size = Pt(12)
            if i == 1:
                textfr.bold = True
    prs.save(outfile)
    print(f'finished inserting regression table for {infile}')


def change_single_info_content(infile, outfile, slide_number, info_text):
    """
    change one slide's information content within the text_frame placeholder
    """
    prs = Presentation(infile)
    slide = prs.slides[slide_number]
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                cur_text = run.text
                new_text = cur_text.replace('INFO', info_text)
                run.text = new_text
    prs.save(outfile)
    print(f'finished changing INFO content for {infile}')

def change_plate_nbr(infile, outfile, info_text):
    """
    change plate number on first slide
    """
    prs = Presentation(infile)
    slide = prs.slides[0]
    for shape in slide.shapes:
        text_frame = shape.text_frame
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                cur_text = run.text
                new_text = cur_text.replace('Plate', info_text)
                run.text = new_text
    prs.save(outfile)
    print(f'finished changing PLATE # content for {infile}')
